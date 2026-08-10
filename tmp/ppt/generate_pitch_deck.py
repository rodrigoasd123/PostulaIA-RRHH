"""Genera entregables/PostulaIA_Pitch.pptx: 5 diapositivas de presentación.

Usa python-pptx. No depende de plantillas externas: construye cada
diapositiva con formas y cuadros de texto para mantener el estilo
consistente con la identidad visual de PostulaIA (azules/navy).
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "entregables" / "PostulaIA_Pitch.pptx"
OUTPUT.parent.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x13, 0x23, 0x3F)
BLUE = RGBColor(0x17, 0x69, 0xE0)
PALE = RGBColor(0xED, 0xF4, 0xFF)
INK = RGBColor(0x17, 0x20, 0x33)
MUTED = RGBColor(0x5C, 0x66, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = align
        para.line_spacing = line_spacing
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size, color,
                 font="Calibri", space_after=10, bold_lead=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"•  {item}"
        para.space_after = Pt(space_after)
        para.line_spacing = 1.15
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_kicker(slide, text):
    add_rect(slide, Inches(0.7), Inches(0.55), Inches(0.5), Inches(0.08), BLUE)
    add_text(slide, Inches(0.7), Inches(0.72), Inches(8), Inches(0.4),
              text.upper(), 13, BLUE, bold=True)


def add_footer(slide, page_num):
    add_text(slide, Inches(0.7), Inches(7.08), Inches(8), Inches(0.3),
              "PostulaIA — Agente de Postulación", 9, MUTED)
    add_text(slide, Inches(12.1), Inches(7.08), Inches(0.6), Inches(0.3),
              str(page_num), 9, MUTED, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide 1 — Portada
# ---------------------------------------------------------------------------
slide = add_slide()
set_background(slide, NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, BLUE)
add_text(slide, Inches(0.9), Inches(2.55), Inches(9), Inches(0.5),
          "PROYECTO DE ANÁLISIS DE CONVOCATORIAS LABORALES", 14, RGBColor(0x9D, 0xC0, 0xFB), bold=True)
add_text(slide, Inches(0.9), Inches(3.0), Inches(11), Inches(1.6),
          "PostulaIA", 60, WHITE, bold=True)
add_text(slide, Inches(0.9), Inches(4.15), Inches(10), Inches(0.7),
          "Agente que analiza convocatorias en PDF y responde preguntas con evidencia citada.",
          18, RGBColor(0xC9, 0xD6, 0xEC))
add_rect(slide, Inches(0.9), Inches(5.15), Inches(1.6), Inches(0.03), BLUE)
add_text(slide, Inches(0.9), Inches(5.35), Inches(10), Inches(0.4),
          "Equipo PostulaIA", 13, RGBColor(0x9D, 0xC0, 0xFB), bold=True)
add_text(slide, Inches(0.9), Inches(5.72), Inches(10), Inches(0.5),
          "Miguel Champi  ·  Rodrigo A.", 15, WHITE)
add_text(slide, Inches(0.9), Inches(6.9), Inches(6), Inches(0.35),
          "MVP local · sin API keys · sin costo", 11, MUTED)

# ---------------------------------------------------------------------------
# Slide 2 — Problema
# ---------------------------------------------------------------------------
slide = add_slide()
set_background(slide, WHITE)
add_kicker(slide, "El problema")
add_text(slide, Inches(0.7), Inches(1.05), Inches(11), Inches(0.9),
          "Leer convocatorias es lento y fácil de malinterpretar", 30, NAVY, bold=True)
add_bullets(slide, Inches(0.7), Inches(2.15), Inches(6.9), Inches(4.2), [
    "Las convocatorias laborales suelen tener 5 a 15 páginas con requisitos, "
    "cronogramas, condiciones y exclusiones dispersos en el texto.",
    "Un postulante puede pasar por alto una fecha límite, un documento "
    "obligatorio o una causal de descalificación escrita en letra pequeña.",
    "Revisar manualmente cada convocatoria toma tiempo y no siempre queda "
    "claro en qué página se sustenta cada requisito.",
    "Las herramientas de IA genéricas responden sin citar la fuente exacta, "
    "lo que dificulta verificar si la respuesta es correcta.",
], 14.5, INK, space_after=14)
card = add_rect(slide, Inches(8.0), Inches(2.15), Inches(4.6), Inches(4.2), PALE)
add_text(slide, Inches(8.35), Inches(2.45), Inches(4.0), Inches(0.4),
          "COSTO DE UN ERROR", 12, BLUE, bold=True)
add_text(slide, Inches(8.35), Inches(2.95), Inches(4.0), Inches(3.2),
          "Perder una postulación por no presentar un documento a tiempo, o "
          "por no cumplir un requisito excluyente que pasó desapercibido "
          "en el texto.",
          15, INK, line_spacing=1.2)
add_footer(slide, 2)

# ---------------------------------------------------------------------------
# Slide 3 — Solución y funcionalidades
# ---------------------------------------------------------------------------
slide = add_slide()
set_background(slide, WHITE)
add_kicker(slide, "La solución")
add_text(slide, Inches(0.7), Inches(1.05), Inches(11.5), Inches(0.9),
          "Un agente local que lee, analiza y responde con evidencia", 28, NAVY, bold=True)
add_text(slide, Inches(0.7), Inches(1.85), Inches(11.5), Inches(0.5),
          "PostulaIA procesa el PDF localmente y estructura la información clave del proceso.",
          14, MUTED)

features = [
    ("Extracción por página", "Lee el PDF y conserva el número de página de cada fragmento de texto."),
    ("Análisis estructurado", "Detecta requisitos, fechas, condiciones laborales y exclusiones."),
    ("Alertas de lectura", "Resalta causales de descalificación y plazos críticos."),
    ("Preguntas con evidencia", "Responde consultas citando la página exacta, ej. [p. 2]."),
    ("Historial local", "Guarda preguntas y respuestas en una base SQLite local."),
    ("IA opcional con Ollama", "Mejora la redacción de respuestas sin enviar datos a la nube."),
]
cols = 3
card_w = Inches(3.65)
card_h = Inches(2.15)
gap_x = Inches(0.25)
gap_y = Inches(0.25)
start_x = Inches(0.7)
start_y = Inches(2.55)
for i, (title, desc) in enumerate(features):
    col = i % cols
    row = i // cols
    x = Emu(start_x + col * (card_w + gap_x))
    y = Emu(start_y + row * (card_h + gap_y))
    add_rect(slide, x, y, card_w, card_h, PALE)
    add_rect(slide, x, y, Inches(0.08), card_h, BLUE)
    add_text(slide, Emu(x + Inches(0.28)), Emu(y + Inches(0.2)), Emu(card_w - Inches(0.5)), Inches(0.5),
              title, 14.5, NAVY, bold=True)
    add_text(slide, Emu(x + Inches(0.28)), Emu(y + Inches(0.75)), Emu(card_w - Inches(0.5)), Emu(card_h - Inches(0.95)),
              desc, 11.5, INK, line_spacing=1.15)
add_footer(slide, 3)

# ---------------------------------------------------------------------------
# Slide 4 — Arquitectura
# ---------------------------------------------------------------------------
slide = add_slide()
set_background(slide, NAVY)
add_kicker(slide, "Arquitectura")
add_text(slide, Inches(0.7), Inches(1.05), Inches(11.5), Inches(0.9),
          "Del PDF a una respuesta verificable", 28, WHITE, bold=True)

pipeline = ["PDF", "Extracción\npor página", "Análisis\n(requisitos y alertas)", "Índice\nléxico", "Respuesta con\nevidencia citada"]
n = len(pipeline)
box_w = Inches(2.05)
box_h = Inches(1.15)
gap = Inches(0.28)
total_w = n * box_w + (n - 1) * gap
start_x = Emu((SLIDE_W - total_w) // 2)
y = Inches(2.75)
for i, label in enumerate(pipeline):
    x = Emu(start_x + i * (box_w + gap))
    color = BLUE if i in (0,) else RGBColor(0x1E, 0x3A, 0x63)
    box = add_rect(slide, x, y, box_w, box_h, color)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    lines = label.split("\n")
    for j, line in enumerate(lines):
        para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = WHITE
    if i < n - 1:
        arrow_x = Emu(x + box_w)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, Emu(y + box_h // 2 - Pt(7)), gap, Pt(14))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x9D, 0xC0, 0xFB)
        arrow.line.fill.background()
        arrow.shadow.inherit = False

opt_y = Emu(y + box_h + Inches(0.5))
opt_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, opt_y, total_w, Inches(0.85))
opt_box.fill.solid()
opt_box.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x63)
opt_box.line.color.rgb = BLUE
opt_box.line.width = Pt(1)
opt_box.shadow.inherit = False
tf = opt_box.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p1 = tf.paragraphs[0]
p1.text = "Ollama local (opcional) mejora la redacción de la respuesta final"
p1.alignment = PP_ALIGN.CENTER
for run in p1.runs:
    run.font.size = Pt(13.5)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x9D, 0xC0, 0xFB)

add_text(slide, Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.5),
          "Todo el procesamiento ocurre en la máquina local. El PDF no se envía a servicios externos.",
          12.5, MUTED)
add_footer(slide, 4)

# ---------------------------------------------------------------------------
# Slide 5 — Resultados y próximos pasos
# ---------------------------------------------------------------------------
slide = add_slide()
set_background(slide, WHITE)
add_kicker(slide, "Estado actual")
add_text(slide, Inches(0.7), Inches(1.05), Inches(11.5), Inches(0.9),
          "Resultados y próximos pasos", 30, NAVY, bold=True)

left_w = Inches(5.6)
add_rect(slide, Inches(0.7), Inches(2.1), left_w, Inches(4.5), PALE)
add_text(slide, Inches(1.0), Inches(2.35), Inches(5.0), Inches(0.4),
          "RESULTADOS", 13, BLUE, bold=True)
add_bullets(slide, Inches(1.0), Inches(2.85), Inches(5.0), Inches(3.6), [
    "MVP funcional en Streamlit, 100% local y sin API keys.",
    "Extracción y análisis automático de requisitos, fechas, "
    "condiciones y exclusiones desde PDF.",
    "Respuestas a preguntas con cita de página como evidencia verificable.",
    "3 convocatorias de ejemplo (TI, administrativo y salud) para "
    "pruebas y demostración.",
    "Suite de pruebas automatizadas (pytest) para el flujo de análisis.",
], 13.5, INK, space_after=12)

right_x = Inches(6.6)
right_w = Inches(6.0)
add_rect(slide, right_x, Inches(2.1), right_w, Inches(4.5), NAVY)
add_text(slide, Emu(right_x + Inches(0.3)), Inches(2.35), Inches(5.4), Inches(0.4),
          "PRÓXIMOS PASOS", 13, RGBColor(0x9D, 0xC0, 0xFB), bold=True)
add_bullets(slide, Emu(right_x + Inches(0.3)), Inches(2.85), Inches(5.4), Inches(3.6), [
    "Soporte OCR para convocatorias escaneadas sin texto seleccionable.",
    "Comparar varias convocatorias en paralelo dentro de una sola sesión.",
    "Exportar el resumen del análisis a PDF o Word.",
    "Empaquetar el modo Ollama con un instalador guiado para el usuario final.",
    "Recolectar retroalimentación de postulantes reales para priorizar mejoras.",
], 13.5, WHITE, space_after=12)

add_footer(slide, 5)

prs.save(str(OUTPUT))
print(OUTPUT)
